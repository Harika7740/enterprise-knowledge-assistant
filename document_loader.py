# document_loader.py - Better DOCX extraction
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image, ImageEnhance   # <-- Added ImageEnhance here
import pytesseract
import io
import streamlit as st

def load_document(file):     
    file_name = file.name.lower()     
    ext = file_name.split('.')[-1]     
    metadata = {"file_name": file.name, "file_type": ext}

    text = ""     
    df = None

    try:         
        file.seek(0)
        if ext == "pdf":             
            reader = PdfReader(file)             
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)             
            metadata["pages"] = len(reader.pages)             
            return text.strip(), metadata, None
        elif ext == "docx":             
            # Better DOCX extraction - include tables, headers, footers             
            doc = Document(file)             
            full_text = []             
            # Paragraphs             
            for para in doc.paragraphs:                 
                if para.text.strip():                     
                    full_text.append(para.text.strip())             
            # Tables             
            for table in doc.tables:                 
                for row in table.rows:                     
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)                     
                    if row_text.strip():                         
                        full_text.append(row_text)             
            # Headers and footers (if any)             
            for section in doc.sections:                 
                header = section.header                 
                for para in header.paragraphs:                     
                    if para.text.strip():                         
                        full_text.append("[Header] " + para.text.strip())                 
                footer = section.footer                 
                for para in footer.paragraphs:                     
                    if para.text.strip():                         
                        full_text.append("[Footer] " + para.text.strip())             
            text = "\n\n".join(full_text)             
            return text.strip() or "No text extracted from DOCX", metadata, None
        elif ext == "pptx":             
            prs = Presentation(file)             
            text = "\n\n".join(shape.text.strip() for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())             
            metadata["slides"] = len(prs.slides)             
            return text.strip(), metadata, None
        elif ext in ["xlsx", "xls"]:             
            df = pd.read_excel(file, engine="openpyxl")             
            text = df.to_string(index=False)             
            metadata["rows"] = df.shape[0]             
            metadata["columns"] = df.shape[1]             
            return text.strip(), metadata, df
        elif ext == "csv":             
            df = pd.read_csv(file)             
            text = df.to_string(index=False)             
            metadata["rows"] = df.shape[0]             
            metadata["columns"] = df.shape[1]             
            return text.strip(), metadata, df
        elif ext == "txt":             
            text = file.read().decode("utf-8", errors="replace")             
            return text.strip(), metadata, None
        elif ext in ["jpg", "jpeg", "png"]:             
            text, metadata = load_image_advanced(file, metadata)             
            return text.strip(), metadata, None
        else:             
            return "", metadata, None
    except Exception as e:         
        st.error(f"Error processing {file.name}: {str(e)}")         
        return f"[ERROR] {str(e)}", metadata, None

# Advanced image OCR (FIXED - only this function changed)
def load_image_advanced(file, metadata):     
    try:         
        file.seek(0)         
        image = Image.open(io.BytesIO(file.read()))         
        image = image.convert('L')         
        width, height = image.size         
        image = image.resize((width * 2, height * 2), Image.LANCZOS)         
        enhancer = ImageEnhance.Contrast(image)   # Fixed: ImageEnhance.Contrast
        image = enhancer.enhance(2.5)         
        enhancer = ImageEnhance.Sharpness(image)  # Fixed: ImageEnhance.Sharpness
        image = enhancer.enhance(1.5)         
        text = pytesseract.image_to_string(image, lang='eng', config='--psm 6')         
        metadata["ocr"] = True         
        return text, metadata     
    except Exception as e:         
        return f"OCR failed: {str(e)}", metadata